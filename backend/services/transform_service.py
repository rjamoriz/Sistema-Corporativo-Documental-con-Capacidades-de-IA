"""
Servicio de Transformación de Documentos
Maneja OCR, conversión de formatos y normalización de contenido
"""
import os
import tempfile
from pathlib import Path
from typing import Dict, Optional

import pytesseract
from pdf2image import convert_from_bytes
from PIL import Image
# import textract  # Commented out - package has dependency issues
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl

from core.logging_config import logger
from models.database_models import Document, DocumentStatus


class TransformService:
    """Servicio para transformación de documentos"""
    
    def __init__(self):
        # Configurar Tesseract
        self.tesseract_config = r'--oem 3 --psm 6'
        # Idiomas soportados: español, inglés, francés, portugués, catalán, euskera, gallego
        self.tesseract_languages = 'spa+eng+fra+por+cat+eus+glg'
    
    async def transform_document(self, document: Document, content: bytes) -> Dict:
        """
        Transforma un documento a texto plano
        
        Args:
            document: Documento a transformar
            content: Contenido binario del documento
            
        Returns:
            Dict: Resultado con texto extraído y metadata
        """
        try:
            mime_type = document.mime_type
            
            # Determinar método de extracción según tipo MIME
            if mime_type == 'application/pdf':
                result = await self._extract_from_pdf(content)
            elif mime_type in ['image/jpeg', 'image/png', 'image/tiff', 'image/bmp']:
                result = await self._extract_from_image(content)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                result = await self._extract_from_docx(content)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                result = await self._extract_from_pptx(content)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                result = await self._extract_from_xlsx(content)
            elif mime_type in ['text/plain', 'text/csv', 'text/html', 'application/json', 'application/xml']:
                result = await self._extract_from_text(content, mime_type)
            else:
                # Intentar con textract como último recurso
                result = await self._extract_with_textract(content, document.filename)
            
            logger.info(f"Document transformed: {document.id}, extracted {len(result['text'])} characters")
            return result
            
        except Exception as e:
            logger.error(f"Error transforming document {document.id}: {e}", exc_info=True)
            return {
                "text": "",
                "page_count": 0,
                "has_images": False,
                "error": str(e)
            }
    
    async def _extract_from_pdf(self, content: bytes) -> Dict:
        """Extrae texto de PDF con OCR si es necesario"""
        text_parts = []
        page_count = 0
        has_images = False
        
        try:
            # Primero intentar extracción directa de texto
            import pdfplumber
            from io import BytesIO
            
            with pdfplumber.open(BytesIO(content)) as pdf:
                page_count = len(pdf.pages)
                
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_parts.append(page_text)
                    else:
                        # Si no hay texto, es probable que sea imagen escaneada
                        has_images = True
            
            # Si no se extrajo suficiente texto, aplicar OCR
            if len(' '.join(text_parts).strip()) < 100:
                logger.info("Insufficient text extracted, applying OCR")
                ocr_text = await self._apply_ocr_to_pdf(content)
                text_parts.append(ocr_text)
                has_images = True
            
            return {
                "text": '\n\n'.join(text_parts),
                "page_count": page_count,
                "has_images": has_images,
                "method": "pdfplumber + OCR" if has_images else "pdfplumber"
            }
            
        except Exception as e:
            logger.warning(f"pdfplumber failed, falling back to OCR: {e}")
            # Si falla, aplicar OCR completo
            ocr_text = await self._apply_ocr_to_pdf(content)
            return {
                "text": ocr_text,
                "page_count": page_count or 1,
                "has_images": True,
                "method": "OCR only"
            }
    
    async def _apply_ocr_to_pdf(self, content: bytes) -> str:
        """Aplica OCR a un PDF convirtiéndolo a imágenes"""
        try:
            # Convertir PDF a imágenes
            images = convert_from_bytes(content, dpi=300)
            
            text_parts = []
            for i, image in enumerate(images):
                logger.debug(f"Processing page {i+1}/{len(images)} with OCR")
                page_text = pytesseract.image_to_string(
                    image,
                    lang=self.tesseract_languages,
                    config=self.tesseract_config
                )
                text_parts.append(page_text)
            
            return '\n\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"OCR failed: {e}", exc_info=True)
            return ""
    
    async def _extract_from_image(self, content: bytes) -> Dict:
        """Extrae texto de imagen usando OCR"""
        try:
            from io import BytesIO
            image = Image.open(BytesIO(content))
            
            # Aplicar OCR
            text = pytesseract.image_to_string(
                image,
                lang=self.tesseract_languages,
                config=self.tesseract_config
            )
            
            return {
                "text": text,
                "page_count": 1,
                "has_images": True,
                "method": "Tesseract OCR",
                "image_size": f"{image.width}x{image.height}"
            }
            
        except Exception as e:
            logger.error(f"Image OCR failed: {e}", exc_info=True)
            return {"text": "", "page_count": 1, "has_images": True, "error": str(e)}
    
    async def _extract_from_docx(self, content: bytes) -> Dict:
        """Extrae texto de DOCX"""
        try:
            from io import BytesIO
            doc = DocxDocument(BytesIO(content))
            
            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extraer texto de tablas
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text)
            
            return {
                "text": '\n\n'.join(text_parts),
                "page_count": len(doc.sections),
                "has_images": len(doc.inline_shapes) > 0,
                "method": "python-docx",
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables)
            }
            
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}", exc_info=True)
            return {"text": "", "page_count": 0, "has_images": False, "error": str(e)}
    
    async def _extract_from_pptx(self, content: bytes) -> Dict:
        """Extrae texto de PPTX"""
        try:
            from io import BytesIO
            prs = Presentation(BytesIO(content))
            
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            
            return {
                "text": '\n\n'.join(text_parts),
                "page_count": len(prs.slides),
                "has_images": True,  # Las presentaciones suelen tener imágenes
                "method": "python-pptx",
                "slide_count": len(prs.slides)
            }
            
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}", exc_info=True)
            return {"text": "", "page_count": 0, "has_images": False, "error": str(e)}
    
    async def _extract_from_xlsx(self, content: bytes) -> Dict:
        """Extrae texto de XLSX"""
        try:
            from io import BytesIO
            wb = openpyxl.load_workbook(BytesIO(content), data_only=True)
            
            text_parts = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"=== {sheet_name} ===")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                    if row_text.strip():
                        text_parts.append(row_text)
            
            return {
                "text": '\n'.join(text_parts),
                "page_count": len(wb.sheetnames),
                "has_images": False,
                "method": "openpyxl",
                "sheet_count": len(wb.sheetnames)
            }
            
        except Exception as e:
            logger.error(f"XLSX extraction failed: {e}", exc_info=True)
            return {"text": "", "page_count": 0, "has_images": False, "error": str(e)}
    
    async def _extract_from_text(self, content: bytes, mime_type: str) -> Dict:
        """Extrae texto de archivos de texto plano"""
        try:
            # Intentar decodificar con diferentes encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            text = None
            
            for encoding in encodings:
                try:
                    text = content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if text is None:
                raise ValueError("Could not decode text with supported encodings")
            
            return {
                "text": text,
                "page_count": 1,
                "has_images": False,
                "method": "direct decoding",
                "mime_type": mime_type
            }
            
        except Exception as e:
            logger.error(f"Text extraction failed: {e}", exc_info=True)
            return {"text": "", "page_count": 1, "has_images": False, "error": str(e)}
    
    async def _extract_with_textract(self, content: bytes, filename: str) -> Dict:
        """Extrae texto usando textract como último recurso"""
        # Textract disabled due to dependency issues
        logger.warning(f"Textract extraction not available for {filename}")
        return {
            "text": "",
            "page_count": 0,
            "has_images": False,
            "method": "textract_unavailable",
            "error": "Textract library has dependency conflicts"
        }
        
        # Original code commented out:
        # try:
        #     # textract requiere archivo temporal
        #     with tempfile.NamedTemporaryFile(delete=False, suffix=Path(filename).suffix) as tmp_file:
        #         tmp_file.write(content)
        #         tmp_path = tmp_file.name
        #     
        #     try:
        #         text = textract.process(tmp_path).decode('utf-8')
        #         return {
        #             "text": text,
        #             "page_count": 1,
        #             "has_images": False,
        #             "method": "textract"
        #         }
        #     finally:
        #         os.unlink(tmp_path)
        #         
        # except Exception as e:
        #     logger.error(f"Textract extraction failed: {e}", exc_info=True)
        #     return {"text": "", "page_count": 1, "has_images": False, "error": str(e)}


# Instancia singleton del servicio
transform_service = TransformService()
